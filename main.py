from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_chroma import Chroma
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from dotenv import load_dotenv
import os
import shutil
from typing import Optional
from pydantic import BaseModel
import uvicorn
from typing import List
from langchain.docstore.document import Document  # Import the Document class

# Load environment variables
load_dotenv()

app = FastAPI()

# Allow cross-origin requests
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize retriever state
retriever = None

# Load persisted Chroma vector store if it exists
def initialize_vectorstore():
    global retriever
    persist_directory = "./chroma_db"
    if os.path.exists(persist_directory):
        try:
            embedding = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            vectorstore = Chroma(persist_directory=persist_directory, embedding_function=embedding)
            retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 10})
            print("Retriever initialized from persisted vector store.")
        except Exception as e:
            print(f"Error loading vector store: {e}")
    else:
        print("Persisted vector store not found.")

initialize_vectorstore()

class QueryRequest(BaseModel):
    input_query: str

@app.post("/upload-files/")
async def upload_files(files: List[UploadFile]):
    processed_files = []  # Accumulate document chunks across all files
    try:
        for file in files:
            temp_file_path = f"./temp_{file.filename}"
            documents = []  # Initialize documents for each file, ensuring it exists in all cases
            try:
                # Save file temporarily
                with open(temp_file_path, "wb") as f:
                    shutil.copyfileobj(file.file, f)

                # Determine file type and load accordingly
                if file.filename.endswith(".pdf"):
                    from langchain_community.document_loaders import PyPDFLoader
                    loader = PyPDFLoader(temp_file_path)
                    data = loader.load()
                    if isinstance(data, list):  # Check if data is a list of Documents
                        documents = data
                    else:  # If not, handle it as a dictionary
                        documents = [Document(page_content=str(doc["page_content"])) for doc in data]
                elif file.filename.endswith(".txt"):
                    from langchain_community.document_loaders import TextLoader
                    loader = TextLoader(temp_file_path)
                    data = loader.load()

                    # Handle data based on its type
                    if isinstance(data, dict):  # If data is a dictionary
                        documents = [Document(page_content=str(value)) for key, value in data.items()]
                    else:  # If data is a list of strings
                        documents = [Document(page_content=str(doc)) for doc in data]
                elif file.filename.endswith(".docx"):
                    from docx import Document as DocxDocument
                    document = DocxDocument(temp_file_path)
                    text = "\n".join([paragraph.text for paragraph in document.paragraphs])
                    documents = [Document(page_content=str(text))]
                elif file.filename.endswith(".ppt"):
                    from pptx import Presentation
                    presentation = Presentation(temp_file_path)
                    text = []
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text.append(shape.text)
                    documents = [Document(page_content=str("\n".join(text)))]
                else:
                    raise ValueError(f"Unsupported file type: {file.filename}")

                # Split documents into chunks
                if documents:
                    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000)
                    docs = text_splitter.split_documents(documents)  # Now docs will be properly defined
                    processed_files.extend(docs)
                else:
                    raise ValueError(f"No data found in the file: {file.filename}")

            finally:
                # Ensure temporary file is removed
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)

        # Create Chroma DB vector store
        if processed_files:
            embedding = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            vectorstore = Chroma.from_documents(
                documents=processed_files,
                embedding=embedding,
                persist_directory="./chroma_db",  # Persistence is automatic
            )
            global retriever
            retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 10})
            return {"message": "All files processed successfully."}
        else:
            return JSONResponse(status_code=400, content={"message": "No valid data in uploaded files."})

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"message": f"Error processing files: {str(e)}"}
        )

@app.post("/query/")
async def query_model(request: QueryRequest):
    try:
        if retriever is None:
            return JSONResponse(status_code=400, content={"message": "Retriever not initialized. Upload a file first."})

        # System prompt for the model
        system_prompt = (
            "You are an assistant for question-answering tasks. "
            "Use the following pieces of retrieved context to answer "
            "the question. If you don't know the answer, say that you "
            "don't know. Use three sentences maximum and keep the "
            "answer concise.\n\n{context}"
        )

        # Create the prompt template
        prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system_prompt),
                ("human", "{input}"),
            ]
        )

        # Set up the question-answer chain
        question_answer_chain = create_stuff_documents_chain(
            ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0),
            prompt
        )
        rag_chain = create_retrieval_chain(retriever, question_answer_chain)

        # Get the response from the RAG chain
        response = rag_chain.invoke({"input": request.input_query})
        return {"answer": response['answer']}

    except Exception as e:
        return JSONResponse(status_code=500, content={"message": f"Error processing query: {str(e)}"})

@app.get("/health/")
def health_check():
    return {"status": "ok"}

if __name__ == '__main__':
    uvicorn.run(app, host='0.0.0.0', port=8000)